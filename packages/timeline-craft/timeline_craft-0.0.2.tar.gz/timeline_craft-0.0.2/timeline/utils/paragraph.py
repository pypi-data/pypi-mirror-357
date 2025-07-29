from pathlib import Path
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT


def amend_font(placeholder, font_family, font_size, bold):
    try:
        text_frame = placeholder.text_frame
        text_frame.fit_text(font_family=font_family, bold=bold, italic=False, max_size=font_size)
    except TypeError:
        amend_font(placeholder, font_family, font_size - 1, bold)


def add_paragraph(
    placeholder,
    text: Path,
    font_size,
    font_family,
    font_color,
):
    """
    Add text to a placeholder; iteratively reduce the size until fits
    """

    text_frame = placeholder.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    # Clear existing paragraphs
    # while len(text_frame.paragraphs) > 0:
    #     text_frame._element.remove(text_frame.paragraphs[0]._element)

    # # Add a new paragraph
    # paragraph = text_frame.add_paragraph()

    paragraph.font.name = font_family
    paragraph.text = text
    paragraph.level = 0
    paragraph.font.color.rgb = font_color
    paragraph.space_before = 0
    paragraph.space_after = 0
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    try:
        text_frame.fit_text(font_family=font_family, bold=False, italic=False, max_size=font_size)
    except TypeError:
        add_paragraph(
            placeholder,
            text,
            font_size - 1,
            font_family=font_family,
            font_color=font_color,
        )
