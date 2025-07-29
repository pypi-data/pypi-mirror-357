import os
from typing import Union
from pptx import Presentation
from datamax.parser.base import BaseLife
from datamax.parser.base import MarkdownOutputVo


class PPtxParser(BaseLife):
    def __init__(self, file_path: Union[str, list]):
        super().__init__()
        self.file_path = file_path

    @staticmethod
    def read_ppt_file(file_path: str):
        try:
            content = ''
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        content += shape.text + '\n'
                    # if shape.shape_type == 13:
                    #     if not os.path.exists("extracted_images"):
                    #         os.makedirs("extracted_images")
                    #     image = shape.image
                    #     image_filename = f'extracted_images/image_{shape.shape_id}.{image.ext}'
                    #     with open(image_filename, 'wb') as img_file:
                    #         img_file.write(image.blob)
                    #     content += ('[' + image_filename + ']')
            return content
        except Exception:
            raise

    def parse(self, file_path: str) -> MarkdownOutputVo:
        try:
            title = os.path.splitext(os.path.basename(file_path))[0]
            content = self.read_ppt_file(file_path=file_path)
            mk_content = content
            lifecycle = self.generate_lifecycle(source_file=file_path, domain="Technology",
                                                usage_purpose="Documentation", life_type="LLM_ORIGIN")
            output_vo = MarkdownOutputVo(title, mk_content)
            output_vo.add_lifecycle(lifecycle)
            return output_vo.to_dict()
        except Exception:
            raise
