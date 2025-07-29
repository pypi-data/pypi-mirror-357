"""
Office document extractor for PowerPoint, Excel, and CSV files.

This module handles extraction of content from Microsoft Office documents
and CSV files, providing structured analysis of presentations, spreadsheets,
and data files.
"""

import csv
import logging
import os
from typing import Any, Optional

from .base_extractor import BaseExtractor
from .models import ExtractionResult

# Enhanced encoding support for international files
try:
    from mark_mate.utils.encoding_utils import safe_read_text_file
    encoding_utils_available = True
except ImportError:
    safe_read_text_file = None  # type: ignore[assignment]
    encoding_utils_available = False

# Optional pandas import
try:
    import pandas as pd
    pandas_available = True
except ImportError:
    pd = None  # type: ignore[assignment]
    pandas_available = False

# Office document imports
try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    Presentation = None  # type: ignore[misc,assignment]
    pptx_available = False

try:
    from openpyxl import load_workbook
    openpyxl_available = True
except ImportError:
    load_workbook = None  # type: ignore[misc,assignment]
    openpyxl_available = False

logger = logging.getLogger(__name__)


class OfficeExtractor(BaseExtractor):
    """Extractor for Office documents (PowerPoint, Excel) and CSV files."""
    
    extractor_name: str
    supported_extensions: list[str]

    def __init__(self, config: Optional[dict[str, Any]] = None):
        super().__init__(config)
        self.extractor_name = "office_extractor"
        self.supported_extensions = [".pptx", ".xlsx", ".xls", ".csv", ".tsv"]

    def can_extract(self, file_path: str) -> bool:
        """Check if this extractor can handle the given file.
        
        Args:
            file_path: Path to the file to check.
            
        Returns:
            True if this is a supported office document, False otherwise.
        """
        ext: str = os.path.splitext(file_path)[1].lower()
        return ext in self.supported_extensions

    def extract_content(self, file_path: str) -> ExtractionResult:
        """Extract content from office documents and CSV files.
        
        Args:
            file_path: Path to the office document to extract.
            
        Returns:
            Dictionary containing extracted content and analysis.
        """
        if not self.can_extract(file_path):
            return self.create_error_result(
                file_path, ValueError("Unsupported file type")
            )

        ext: str = os.path.splitext(file_path)[1].lower()

        try:
            if ext == ".pptx":
                return self._extract_powerpoint(file_path)
            elif ext in [".xlsx", ".xls"]:
                return self._extract_excel(file_path)
            elif ext in [".csv", ".tsv"]:
                return self._extract_csv(file_path)
            else:
                return self.create_error_result(
                    file_path, ValueError(f"Unsupported extension: {ext}")
                )

        except Exception as e:
            logger.error(f"Error extracting office document {file_path}: {e}")
            return self.create_error_result(file_path, e)

    def _extract_powerpoint(self, file_path: str) -> ExtractionResult:
        """Extract content from PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint file.
            
        Returns:
            Dictionary containing extracted presentation content.
        """
        if not pptx_available:
            error_msg: str = "python-pptx not available for PowerPoint extraction"
            logger.warning(error_msg)
            return self.create_error_result(file_path, ImportError(error_msg))

        try:
            presentation = Presentation(file_path)
            slides_content: list[dict[str, Any]] = []
            total_text: str = ""

            for i, slide in enumerate(presentation.slides, 1):
                slide_data: dict[str, Any] = {
                    "slide_number": i,
                    "title": "",
                    "content": "",
                    "notes": "",
                }

                # Extract title and content from shapes
                slide_text_parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text_parts.append(text)

                        # Try to identify title (usually first text box or larger font)
                        if not slide_data["title"] and len(text) < 100:
                            slide_data["title"] = text
                        else:
                            slide_data["content"] += text + "\\n"

                # Extract speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_data["notes"] = (
                        slide.notes_slide.notes_text_frame.text.strip()
                    )

                # Combine all text for this slide
                all_slide_text = "\\n".join(slide_text_parts)
                if slide_data["notes"]:
                    all_slide_text += "\\n[SPEAKER NOTES]\\n" + slide_data["notes"]

                slides_content.append(slide_data)
                total_text += f"\\n=== Slide {i} ===\\n" + all_slide_text + "\\n"

            # Create analysis
            analysis = {
                "total_slides": len(presentation.slides),
                "slides_with_notes": sum(1 for s in slides_content if s["notes"]),
                "avg_content_length": sum(len(s["content"]) for s in slides_content)
                / len(slides_content)
                if slides_content
                else 0,
                "slides_breakdown": slides_content[
                    :10
                ],  # First 10 slides for detailed analysis
            }

            # Create content summary
            content = f"""POWERPOINT PRESENTATION SUMMARY:
Total Slides: {analysis["total_slides"]}
Slides with Speaker Notes: {analysis["slides_with_notes"]}
Average Content Length: {analysis["avg_content_length"]:.1f} characters

PRESENTATION CONTENT:
{total_text}"""

            return self.create_success_result(file_path, content, analysis)

        except Exception as e:
            logger.error(f"Error extracting PowerPoint {file_path}: {e}")
            return self.create_error_result(file_path, e)

    def _extract_excel(self, file_path: str) -> ExtractionResult:
        """Extract content from Excel workbook.
        
        Args:
            file_path: Path to the Excel file.
            
        Returns:
            Dictionary containing extracted spreadsheet content.
        """
        if not openpyxl_available:
            error_msg: str = "openpyxl not available for Excel extraction"
            logger.warning(error_msg)
            return self.create_error_result(file_path, ImportError(error_msg))

        try:
            # Try to load with pandas first for better data analysis
            if pandas_available:
                try:
                    # Read all sheets
                    excel_data = pd.read_excel(
                        file_path, sheet_name=None, engine="openpyxl"
                    )
                    sheets_analysis: dict[str, Any] = {}
                    total_content: str = ""

                    for sheet_name, df in excel_data.items():
                        # Basic data analysis
                        sheet_analysis = {
                            "name": sheet_name,
                            "shape": df.shape,
                            "columns": list(df.columns),
                            "dtypes": df.dtypes.to_dict(),
                            "null_counts": df.isnull().sum().to_dict(),
                            "sample_data": df.head(3).to_dict() if not df.empty else {},
                        }

                        # Add basic statistics for numeric columns
                        numeric_cols = df.select_dtypes(include=["number"]).columns
                        if len(numeric_cols) > 0:
                            sheet_analysis["numeric_summary"] = (
                                df[numeric_cols].describe().to_dict()
                            )

                        sheets_analysis[sheet_name] = sheet_analysis

                        # Create text representation
                        sheet_text = f"\\n=== Sheet: {sheet_name} ===\\n"
                        sheet_text += (
                            f"Dimensions: {df.shape[0]} rows × {df.shape[1]} columns\\n"
                        )
                        sheet_text += f"Columns: {', '.join(df.columns)}\\n"

                        if not df.empty:
                            sheet_text += "\\nFirst few rows:\\n"
                            sheet_text += df.head(5).to_string() + "\\n"

                        total_content += sheet_text

                    # Overall analysis
                    analysis = {
                        "total_sheets": len(excel_data),
                        "total_rows": sum(df.shape[0] for df in excel_data.values()),
                        "total_columns": sum(df.shape[1] for df in excel_data.values()),
                        "sheets_analysis": sheets_analysis,
                    }

                    content = f"""EXCEL WORKBOOK SUMMARY:
Total Sheets: {analysis["total_sheets"]}
Total Rows: {analysis["total_rows"]}
Total Columns: {analysis["total_columns"]}

WORKBOOK CONTENT:
{total_content}"""

                    return self.create_success_result(file_path, content, analysis)

                except Exception as pandas_error:
                    logger.warning(
                        f"Pandas Excel reading failed for {file_path}, trying openpyxl directly: {pandas_error}"
                    )
            else:
                logger.info(
                    f"Pandas not available, using openpyxl directly for {file_path}"
                )

            # Fallback to openpyxl direct reading
            workbook = load_workbook(file_path, data_only=True)
            sheets_content = []
            total_content = ""

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"\\n=== Sheet: {sheet_name} ===\\n"

                # Extract cell values (limit to reasonable size)
                rows_content = []
                max_rows = min(sheet.max_row, 100)  # Limit to first 100 rows
                max_cols = min(sheet.max_column, 20)  # Limit to first 20 columns

                for row in sheet.iter_rows(
                    min_row=1,
                    max_row=max_rows,
                    min_col=1,
                    max_col=max_cols,
                    values_only=True,
                ):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):  # Skip empty rows
                        rows_content.append(" | ".join(row_values))

                sheet_text += "\\n".join(rows_content[:20])  # First 20 rows
                sheets_content.append(
                    {
                        "name": sheet_name,
                        "rows": sheet.max_row,
                        "columns": sheet.max_column,
                        "content_preview": sheet_text,
                    }
                )
                total_content += sheet_text + "\\n"

            analysis = {
                "total_sheets": len(workbook.sheetnames),
                "sheets": sheets_content,
                "extraction_method": "openpyxl_fallback",
            }

            content = f"""EXCEL WORKBOOK SUMMARY (Basic Extraction):
Total Sheets: {len(workbook.sheetnames)}
Sheet Names: {", ".join(workbook.sheetnames)}

WORKBOOK CONTENT:
{total_content}"""

            return self.create_success_result(file_path, content, analysis)

        except Exception as e:
            logger.error(f"Error extracting Excel {file_path}: {e}")
            return self.create_error_result(file_path, e)

    def _extract_csv(self, file_path: str) -> ExtractionResult:
        """Extract content from CSV/TSV file.
        
        Args:
            file_path: Path to the CSV/TSV file.
            
        Returns:
            Dictionary containing extracted CSV content and analysis.
        """
        try:
            # Detect delimiter
            delimiter: str = "\\t" if file_path.lower().endswith(".tsv") else ","

            # Enhanced encoding detection for international CSV files
            df = None
            used_encoding: Optional[str] = None
            csv_content: Optional[str] = None

            # Try enhanced encoding support first
            if encoding_utils_available and safe_read_text_file is not None:
                csv_content, used_encoding, error_msg = safe_read_text_file(file_path, "csv")
                if csv_content is not None:
                    logger.info(f"Successfully read CSV with enhanced encoding: {used_encoding}")
                else:
                    logger.warning(f"Enhanced encoding failed for {file_path}: {error_msg}")

            # If enhanced encoding worked and pandas is available, use pandas with the detected encoding
            if csv_content is not None and pandas_available and used_encoding:
                try:
                    # Clean the encoding name for pandas (remove error handling info)
                    clean_encoding = used_encoding.split(' (')[0] if ' (' in used_encoding else used_encoding
                    df = pd.read_csv(file_path, delimiter=delimiter, encoding=clean_encoding)
                    logger.info(f"Successfully read CSV with pandas using detected encoding: {clean_encoding}")
                except Exception as e:
                    clean_encoding = used_encoding.split(' (')[0] if ' (' in used_encoding else used_encoding
                    logger.warning(f"Pandas failed with detected encoding {clean_encoding}: {e}")
                    # Fall back to basic encoding attempts
                    fallback_encodings = ["utf-8", "latin-1", "cp1252"]
                    for encoding in fallback_encodings:
                        try:
                            df = pd.read_csv(file_path, delimiter=delimiter, encoding=encoding)
                            used_encoding = encoding
                            logger.info(f"Pandas fallback successful with {encoding}")
                            break
                        except Exception:
                            continue
            
            # If no enhanced encoding or pandas failed, try basic pandas encoding attempts
            elif pandas_available:
                fallback_encodings = ["utf-8", "latin-1", "cp1252", "cp1251", "gb2312", "shift_jis"]
                for encoding in fallback_encodings:
                    try:
                        df = pd.read_csv(file_path, delimiter=delimiter, encoding=encoding)
                        used_encoding = encoding
                        logger.info(f"Basic pandas encoding successful with {encoding}")
                        break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        logger.warning(f"Error reading CSV with pandas and {encoding}: {e}")
                        continue

            if df is not None:
                # Pandas-based analysis
                analysis = {
                    "shape": df.shape,
                    "columns": list(df.columns),
                    "dtypes": df.dtypes.to_dict(),
                    "null_counts": df.isnull().sum().to_dict(),
                    "null_percentage": (df.isnull().sum() / len(df) * 100).to_dict(),
                    "encoding_used": used_encoding,
                    "delimiter": "tab" if delimiter == "\\t" else "comma",
                }

                # Add statistics for numeric columns
                numeric_cols = df.select_dtypes(include=["number"]).columns
                if len(numeric_cols) > 0:
                    analysis["numeric_summary"] = df[numeric_cols].describe().to_dict()

                # Add sample data
                analysis["sample_data"] = {
                    "first_5_rows": df.head(5).to_dict(),
                    "last_5_rows": df.tail(5).to_dict() if len(df) > 5 else {},
                }

                # Create content text
                content = f"""CSV/TSV DATA SUMMARY:
File: {os.path.basename(file_path)}
Dimensions: {df.shape[0]} rows × {df.shape[1]} columns
Encoding: {used_encoding}
Delimiter: {analysis["delimiter"]}

COLUMNS:
{", ".join(df.columns)}

DATA QUALITY:
"""

                # Add data quality information
                for col in df.columns:
                    null_percentage = analysis.get("null_percentage")
                    if null_percentage is not None:
                        null_pct = null_percentage[col]
                        content += f"- {col}: {null_pct:.1f}% missing\\n"

                content += f"\\nFIRST FEW ROWS:\\n{df.head(10).to_string()}\\n"

                if len(df) > 10:
                    content += f"\\nLAST FEW ROWS:\\n{df.tail(5).to_string()}\\n"

                return self.create_success_result(file_path, content, analysis)
            else:
                # Fallback: basic CSV reading
                logger.info(f"Using basic CSV reader for {file_path}")
                rows = []
                headers = None

                # If we have CSV content from enhanced encoding, parse it directly
                if csv_content is not None:
                    try:
                        from io import StringIO
                        reader = csv.reader(StringIO(csv_content), delimiter=delimiter)
                        rows = list(reader)
                        logger.info(f"Successfully parsed CSV content with enhanced encoding: {used_encoding}")
                    except Exception as e:
                        logger.warning(f"Failed to parse CSV content from enhanced encoding: {e}")
                        csv_content = None
                
                # If enhanced encoding didn't work, try basic encoding fallback
                if not rows:
                    fallback_encodings = ["utf-8", "latin-1", "cp1252", "cp1251", "gb2312", "shift_jis", "big5"]
                    for encoding in fallback_encodings:
                        try:
                            with open(file_path, encoding=encoding) as f:
                                reader = csv.reader(f, delimiter=delimiter)
                                rows = list(reader)
                                used_encoding = encoding
                                logger.info(f"Basic CSV reader successful with {encoding}")
                                break
                        except UnicodeDecodeError:
                            continue
                        except Exception as e:
                            logger.warning(f"Error reading CSV with basic reader and {encoding}: {e}")
                            continue

                if not rows:
                    raise ValueError("Could not read CSV file with any encoding")

                headers = rows[0] if rows else []
                data_rows = rows[1:] if len(rows) > 1 else []

                # Basic analysis
                analysis = {
                    "shape": (len(data_rows), len(headers)),
                    "columns": headers,
                    "total_rows": len(data_rows),
                    "encoding_used": used_encoding,
                    "delimiter": repr(delimiter),
                    "analysis_method": "basic_csv_reader",
                }

                # Create content text
                content = f"""CSV/TSV DATA SUMMARY:
File: {os.path.basename(file_path)}
Dimensions: {len(data_rows)} rows × {len(headers)} columns
Encoding: {used_encoding}
Delimiter: {repr(delimiter)}

COLUMNS:
{", ".join(headers)}

FIRST FEW ROWS:
"""
                # Show first 10 rows
                for i, row in enumerate(data_rows[:10]):
                    content += f"Row {i + 1}: {', '.join(row[:5])}{'...' if len(row) > 5 else ''}\n"

                return self.create_success_result(file_path, content, analysis)

        except Exception as e:
            logger.error(f"Error extracting CSV {file_path}: {e}")
            return self.create_error_result(file_path, e)
