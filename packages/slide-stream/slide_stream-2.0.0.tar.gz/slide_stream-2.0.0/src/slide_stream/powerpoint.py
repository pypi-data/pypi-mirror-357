"""PowerPoint parsing functionality for Slide Stream."""

from pathlib import Path
from typing import Any

from pptx import Presentation


def parse_powerpoint(file_path: str | Path) -> list[dict[str, Any]]:
    """Parse PowerPoint file into slide data.
    
    Extracts slide titles, bullet points, and speaker notes from .pptx files.
    
    Args:
        file_path: Path to the PowerPoint file (.pptx)
        
    Returns:
        List of slide dictionaries with 'title', 'content', and 'notes' keys
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Could not open PowerPoint file: {e}") from e
    
    slides = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "title": f"Slide {slide_num}",  # Default title
            "content": [],
            "notes": ""
        }
        
        # Extract slide content
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Check if this might be a title (usually the first text shape or larger font)
                text = shape.text.strip()
                if not text:
                    continue
                
                # If this is the first significant text and looks like a title
                if (slide_data["title"] == f"Slide {slide_num}" and 
                    len(text) < 100 and 
                    '\n' not in text):
                    slide_data["title"] = text
                else:
                    # Extract bullet points or paragraphs
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_data["content"].append(para_text)
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if shape.has_text_frame:
                    notes_text = shape.text.strip()
                    if notes_text and notes_text != "Click to add notes":
                        slide_data["notes"] = notes_text
                        break
        
        # Only add slides that have meaningful content
        if slide_data["content"] or slide_data["notes"]:
            slides.append(slide_data)
    
    return slides


def format_powerpoint_content_for_llm(slide: dict[str, Any]) -> str:
    """Format PowerPoint slide data for LLM processing.
    
    Combines title, content, and speaker notes into a coherent prompt.
    
    Args:
        slide: Slide dictionary with title, content, and notes
        
    Returns:
        Formatted string for LLM input
    """
    parts = [f"Title: {slide['title']}"]
    
    if slide["content"]:
        parts.append("Content:")
        for item in slide["content"]:
            parts.append(f"• {item}")
    
    if slide["notes"]:
        parts.append(f"Speaker Notes: {slide['notes']}")
    
    return "\n".join(parts)