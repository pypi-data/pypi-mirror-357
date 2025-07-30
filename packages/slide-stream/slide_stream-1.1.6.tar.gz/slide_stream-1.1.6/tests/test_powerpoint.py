"""Tests for PowerPoint parsing functionality."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from slide_stream.powerpoint import format_powerpoint_content_for_llm, parse_powerpoint


def create_test_powerpoint() -> Path:
    """Create a test PowerPoint file with sample content."""
    temp_file = Path(tempfile.mktemp(suffix=".pptx"))
    
    # Create a presentation
    prs = Presentation()
    
    # Slide 1: Title and content
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title1 = slide1.shapes.title
    title1.text = "Introduction to AI"
    
    content1 = slide1.placeholders[1]
    content1.text = "What is Artificial Intelligence?\nMachine Learning basics\nReal-world applications"
    
    # Add speaker notes
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "Welcome everyone to this presentation on AI. Start with a brief overview."
    
    # Slide 2: Another slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Machine Learning Types"
    
    content2 = slide2.placeholders[1]
    content2.text = "Supervised Learning\nUnsupervised Learning\nReinforcement Learning"
    
    # Save the presentation
    prs.save(temp_file)
    return temp_file


def test_parse_powerpoint_basic():
    """Test basic PowerPoint parsing."""
    pptx_file = create_test_powerpoint()
    
    try:
        slides = parse_powerpoint(pptx_file)
        
        assert len(slides) == 2
        
        # Check first slide
        assert slides[0]["title"] == "Introduction to AI"
        assert "What is Artificial Intelligence?" in slides[0]["content"]
        assert "Machine Learning basics" in slides[0]["content"]
        assert "Real-world applications" in slides[0]["content"]
        assert "Welcome everyone" in slides[0]["notes"]
        
        # Check second slide
        assert slides[1]["title"] == "Machine Learning Types"
        assert "Supervised Learning" in slides[1]["content"]
        assert "Unsupervised Learning" in slides[1]["content"]
        assert "Reinforcement Learning" in slides[1]["content"]
        
    finally:
        pptx_file.unlink(missing_ok=True)


def test_parse_powerpoint_nonexistent_file():
    """Test parsing non-existent PowerPoint file."""
    with pytest.raises(ValueError, match="Could not open PowerPoint file"):
        parse_powerpoint("nonexistent.pptx")


def test_format_powerpoint_content_for_llm():
    """Test formatting PowerPoint content for LLM input."""
    slide_data = {
        "title": "Test Slide",
        "content": ["Point 1", "Point 2", "Point 3"],
        "notes": "These are speaker notes for the slide."
    }
    
    formatted = format_powerpoint_content_for_llm(slide_data)
    
    assert "Title: Test Slide" in formatted
    assert "Content:" in formatted
    assert "• Point 1" in formatted
    assert "• Point 2" in formatted
    assert "• Point 3" in formatted
    assert "Speaker Notes: These are speaker notes for the slide." in formatted


def test_format_powerpoint_content_no_notes():
    """Test formatting PowerPoint content without speaker notes."""
    slide_data = {
        "title": "Test Slide",
        "content": ["Point 1", "Point 2"],
        "notes": ""
    }
    
    formatted = format_powerpoint_content_for_llm(slide_data)
    
    assert "Title: Test Slide" in formatted
    assert "• Point 1" in formatted
    assert "• Point 2" in formatted
    assert "Speaker Notes:" not in formatted


def test_format_powerpoint_content_no_content():
    """Test formatting PowerPoint with only notes."""
    slide_data = {
        "title": "Notes Only Slide",
        "content": [],
        "notes": "Just speaker notes here."
    }
    
    formatted = format_powerpoint_content_for_llm(slide_data)
    
    assert "Title: Notes Only Slide" in formatted
    assert "Content:" not in formatted
    assert "Speaker Notes: Just speaker notes here." in formatted


def test_empty_powerpoint():
    """Test parsing PowerPoint with no meaningful content."""
    temp_file = Path(tempfile.mktemp(suffix=".pptx"))
    
    try:
        # Create empty presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        prs.save(temp_file)
        
        slides = parse_powerpoint(temp_file)
        assert len(slides) == 0  # No meaningful content
        
    finally:
        temp_file.unlink(missing_ok=True)