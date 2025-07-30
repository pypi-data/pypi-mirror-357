"""Tests for the CLI interface."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from typer.testing import CliRunner

from slide_stream.cli import app


@pytest.fixture
def runner():
    """Create a CLI test runner."""
    return CliRunner()


@pytest.fixture
def sample_markdown():
    """Create sample markdown content."""
    return """# Test Slide

- First point
- Second point

# Another Slide

- More content
- Final point
"""


@pytest.fixture
def sample_powerpoint():
    """Create sample PowerPoint file."""
    temp_file = Path(tempfile.mktemp(suffix=".pptx"))
    
    # Create a presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    
    title = slide.shapes.title
    title.text = "Test PowerPoint Slide"
    
    content = slide.placeholders[1]
    content.text = "First bullet point\nSecond bullet point"
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are test speaker notes."
    
    prs.save(temp_file)
    return temp_file


def test_cli_help(runner):
    """Test CLI help command."""
    result = runner.invoke(app, ["--help"])
    assert result.exit_code == 0
    assert "Create a video from a Markdown (.md) or PowerPoint (.pptx) file" in result.stdout


def test_cli_version(runner):
    """Test CLI version command."""
    result = runner.invoke(app, ["--version"])
    assert result.exit_code == 0
    assert "SlideStream" in result.stdout
    assert "1.1.7" in result.stdout


def test_cli_create_missing_input(runner):
    """Test CLI with missing input file."""
    result = runner.invoke(app, [])
    assert result.exit_code != 0


def test_cli_create_basic(runner, sample_markdown):
    """Test basic CLI create command."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False) as f:
        f.write(sample_markdown)
        f.flush()
        
        # Test with text image source to avoid network calls
        result = runner.invoke(app, [
            "--input", f.name,
            "--output", "test_output.mp4",
            "--image-source", "text"
        ])
        
        # Should not crash during parsing phase
        assert "Parsing Markdown" in result.stdout or result.exit_code == 0
        
        # Clean up
        Path(f.name).unlink(missing_ok=True)
        Path("test_output.mp4").unlink(missing_ok=True)


def test_cli_create_powerpoint(runner, sample_powerpoint):
    """Test CLI create command with PowerPoint file."""
    try:
        # Test with PowerPoint file and text image source to avoid network calls
        result = runner.invoke(app, [
            "--input", str(sample_powerpoint),
            "--output", "test_pptx_output.mp4",
            "--image-source", "text"
        ])
        
        # Should not crash during parsing phase
        assert "Parsing PowerPoint" in result.stdout or result.exit_code == 0
        
        # Clean up
        Path("test_pptx_output.mp4").unlink(missing_ok=True)
        
    finally:
        sample_powerpoint.unlink(missing_ok=True)


def test_cli_unsupported_file_type(runner):
    """Test CLI with unsupported file type."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        f.write("Some text content")
        f.flush()
        
        result = runner.invoke(app, [
            "--input", f.name,
            "--output", "test_output.mp4"
        ])
        
        assert result.exit_code != 0
        assert "Unsupported file type" in result.stdout
        
        # Clean up
        Path(f.name).unlink(missing_ok=True)


def test_cli_nonexistent_file(runner):
    """Test CLI with non-existent input file."""
    result = runner.invoke(app, [
        "--input", "nonexistent.md",
        "--output", "test_output.mp4"
    ])
    
    assert result.exit_code != 0
    assert "Input file not found" in result.stdout