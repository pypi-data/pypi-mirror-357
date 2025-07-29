from pptx import Presentation
import os

class PPTXExtractor:
    def __init__(self, pptx_path, image_output_dir="extracted_images"):
        self.pptx_path = pptx_path
        self.image_output_dir = image_output_dir
        self.prs = Presentation(pptx_path)

        if not os.path.exists(self.image_output_dir):
            os.makedirs(self.image_output_dir)

    def extract_text(self):
        text_data = {}
        for i, slide in enumerate(self.prs.slides, start=1):
            text_data[i] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_data[i].append(text)
        return text_data

    def extract_images(self):
        image_data = {}
        for slide_num, slide in enumerate(self.prs.slides, start=1):
            image_data[slide_num] = []
            for shape_num, shape in enumerate(slide.shapes, start=1):
                if shape.shape_type == 13:  # Picture shape type
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = f"slide{slide_num}_image{shape_num}.{image_ext}"
                    image_path = os.path.join(self.image_output_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    image_data[slide_num].append(image_path)
        return image_data

    def extract_tables(self):
        table_data = {}
        for slide_num, slide in enumerate(self.prs.slides, start=1):
            table_data[slide_num] = []
            for shape_num, shape in enumerate(slide.shapes, start=1):
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_rows.append(row_data)
                    table_data[slide_num].append(table_rows)
        return table_data

    def extract_all(self):
        return {
            "texts": self.extract_text(),
            "images": self.extract_images(),
            "tables": self.extract_tables()
        }


if __name__ == "__main__":
    extractor = PPTXExtractor("your_presentation.pptx")

    data = extractor.extract_all()

    for slide_num in data["texts"].keys():
        print(f"\nSlide {slide_num} Texts:")
        for t in data["texts"][slide_num]:
            print("-", t)

        print(f"Slide {slide_num} Images:")
        for img_path in data["images"].get(slide_num, []):
            print("-", img_path)

        print(f"Slide {slide_num} Tables:")
        for table in data["tables"].get(slide_num, []):
            for row in table:
                print("\t".join(row))
            print("---")
